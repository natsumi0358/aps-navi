import os
import io
import json
import anthropic
import requests
from flask import Flask, render_template, request, jsonify, session, redirect, url_for, send_file
from pathlib import Path
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse
from database import init_db, get_all_companies, get_company, create_company, update_company, delete_company

app = Flask(__name__)
app.secret_key = os.environ.get("FLASK_SECRET_KEY", "nknarts-aps-navi-secret")

# DB初期化（gunicorn起動時にも必ず実行）
init_db()

# システムプロンプトを読み込む
SYSTEM_PROMPT_PATH = Path(__file__).parent.parent / "aps-advisor-ai" / "knowledge" / "system_prompt.md"
try:
    with open(SYSTEM_PROMPT_PATH, encoding="utf-8") as f:
        BASE_SYSTEM_PROMPT = f.read()
except FileNotFoundError:
    BASE_SYSTEM_PROMPT = "あなたはAPS観点で営業アドバイスをするAIアドバイザーです。"

# PPTXテンプレートパス
TEMPLATE_PATH = Path(__file__).parent / "アカウントプラン_テンプレート_260226.pptx"

client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))
MAX_HISTORY = 20

# 会社概要ページを示すキーワード（URLに含まれる場合に追加取得）
COMPANY_INFO_KEYWORDS = [
    "company", "about", "corporate", "profile", "gaisha", "outline",
    "kaisha", "gaiyo", "概要", "会社情報", "企業情報", "会社案内"
]


# ─────────────────────────────────────────
# ヘルパー関数
# ─────────────────────────────────────────

def build_system_prompt(company):
    systems = company.get("systems", [])
    key_persons = company.get("key_persons", [])

    systems_text = "\n".join([
        f"  - {s.get('name','')}（{s.get('type','')}）導入：{s.get('installed','')} / リプレース見込み：{s.get('replace','')}"
        for s in systems
    ]) or "  （未入力）"

    key_persons_text = "\n".join([
        f"  - {p.get('name','')}（{p.get('title','')}）意思決定権：{p.get('decision','')} / {p.get('note','')}"
        for p in key_persons
    ]) or "  （未入力）"

    company_info = f"""
【調査対象会社】{company.get('company_name', '')}
【業種】{company.get('industry', '')}
【従業員数】{company.get('employees', '')}
【売上規模】{company.get('revenue', '')}
【担当営業】{company.get('sales_person', '')}

■ Work6: お客様概要
【事業内容】{company.get('overview', '（未入力）')}
【MVV】{company.get('mvv', '（未入力）')}
【代表者プロフィール】{company.get('president_profile', '（未入力）')}

■ Work7: 経営方針
【中期経営計画】{company.get('mid_term_plan', '（未入力）')}
【IR情報】{company.get('ir_info', '（未入力）')}
【投資エリア・想定ソリューション】{company.get('investment_areas', '（未入力）')}

■ 業界分析
【PEST分析】{company.get('pest', '（未入力）')}
【競合状況】{company.get('competitors', '（未入力）')}

■ 営業情報
【エンドユーザーの課題】{company.get('end_user_issues', '（未入力）')}
【潜在ニーズ】{company.get('latent_needs', '（未入力）')}
【ビッグプレー候補】{company.get('big_play', '（未入力）')}
【パイプライン状況】{company.get('pipeline', '（未入力）')}
【キーパーソン】{key_persons_text}
【システム一覧】{systems_text}

上記の情報をもとに、APS観点でアドバイスをしてください。
"""
    return BASE_SYSTEM_PROMPT + "\n\n---\n" + company_info


def fetch_page_text(url, headers):
    """1ページ分のテキストを取得（失敗時はNone）"""
    try:
        resp = requests.get(url, headers=headers, timeout=10)
        resp.encoding = resp.apparent_encoding
        soup = BeautifulSoup(resp.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True), soup
    except Exception:
        return None, None


def fetch_url_text(url):
    """
    URLのHTMLテキストを取得。
    メインページに加え、会社概要ページも自動検出して取得・結合する。
    """
    headers = {"User-Agent": "Mozilla/5.0 (compatible; APSNavi/1.0)"}
    main_text, soup = fetch_page_text(url, headers)
    if not main_text:
        raise Exception("URLの取得に失敗しました")

    combined_text = main_text[:5000]

    # 会社概要サブページを探す
    if soup:
        base = f"{urlparse(url).scheme}://{urlparse(url).netloc}"
        found_sub = set()
        for a in soup.find_all("a", href=True):
            href = a["href"].lower()
            link_text = a.get_text(strip=True).lower()
            # URL または リンクテキストにキーワードが含まれるリンクを探す
            is_info_page = any(kw in href for kw in COMPANY_INFO_KEYWORDS) or \
                           any(kw in link_text for kw in ["会社概要", "企業情報", "会社情報", "会社案内", "corporate", "about"])
            if is_info_page:
                full_url = urljoin(url, a["href"])
                # 同じドメインのみ
                if urlparse(full_url).netloc == urlparse(url).netloc:
                    found_sub.add(full_url)

        # 最大2ページまで追加取得
        for sub_url in list(found_sub)[:2]:
            if sub_url != url:
                sub_text, _ = fetch_page_text(sub_url, headers)
                if sub_text:
                    combined_text += "\n\n--- 会社概要ページ ---\n" + sub_text[:3000]

    return combined_text


def set_slide_text(slide, placeholder_text, new_text):
    """スライド内の指定テキストを置換"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if placeholder_text in run.text:
                    run.text = run.text.replace(placeholder_text, new_text)


def write_to_textbox(slide, search_keyword, new_text, max_chars=800):
    """キーワードを含むテキストボックスの内容を置き換える"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        full_text = shape.text_frame.text
        if search_keyword in full_text:
            tf = shape.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            if tf.paragraphs:
                if tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = str(new_text)[:max_chars]
                else:
                    tf.paragraphs[0].text = str(new_text)[:max_chars]
            return True
    return False


# ─────────────────────────────────────────
# ルート
# ─────────────────────────────────────────

@app.route("/")
def index():
    companies = get_all_companies()
    return render_template("index.html", companies=companies)


@app.route("/company/new", methods=["GET", "POST"])
def company_new():
    if request.method == "POST":
        data = _parse_form(request.form)
        company_id = create_company(data)
        action = request.form.get("action_type", "save")
        if action == "save_ppt":
            return redirect(url_for("download_ppt", company_id=company_id))
        return redirect(url_for("company_advisor", company_id=company_id))
    return render_template("company_form.html", company=None, mode="new")


@app.route("/company/<int:company_id>/edit", methods=["GET", "POST"])
def company_edit(company_id):
    company = get_company(company_id)
    if not company:
        return redirect(url_for("index"))
    if request.method == "POST":
        data = _parse_form(request.form)
        update_company(company_id, data)
        action = request.form.get("action_type", "save")
        if action == "save_ppt":
            return redirect(url_for("download_ppt", company_id=company_id))
        return redirect(url_for("company_advisor", company_id=company_id))
    return render_template("company_form.html", company=company, mode="edit")


@app.route("/company/<int:company_id>/advisor")
def company_advisor(company_id):
    company = get_company(company_id)
    if not company:
        return redirect(url_for("index"))
    session_key = f"history_{company_id}"
    session[session_key] = []
    return render_template("advisor.html", company=company)


@app.route("/company/<int:company_id>/chat", methods=["POST"])
def company_chat(company_id):
    company = get_company(company_id)
    if not company:
        return jsonify({"error": "会社情報が見つかりません"}), 404

    data = request.get_json()
    user_message = data.get("message", "").strip()
    if not user_message:
        return jsonify({"error": "メッセージを入力してください"}), 400

    session_key = f"history_{company_id}"
    history = session.get(session_key, [])
    history.append({"role": "user", "content": user_message})
    if len(history) > MAX_HISTORY:
        history = history[-MAX_HISTORY:]

    try:
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=2048,
            system=build_system_prompt(company),
            messages=history,
        )
        assistant_message = response.content[0].text
        history.append({"role": "assistant", "content": assistant_message})
        session[session_key] = history
        return jsonify({"reply": assistant_message})
    except Exception as e:
        return jsonify({"error": f"エラーが発生しました: {str(e)}"}), 500


@app.route("/company/<int:company_id>/analyze", methods=["POST"])
def company_analyze(company_id):
    company = get_company(company_id)
    if not company:
        return jsonify({"error": "会社情報が見つかりません"}), 404

    first_message = f"{company.get('company_name', '')}についてAPS観点で総合的にアドバイスしてください。特に「ビッグプレー候補」「潜在ニーズの深掘り」「パイプラインの充実度」の観点でお願いします。"
    history = [{"role": "user", "content": first_message}]

    try:
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=2048,
            system=build_system_prompt(company),
            messages=history,
        )
        assistant_message = response.content[0].text
        history.append({"role": "assistant", "content": assistant_message})
        session[f"history_{company_id}"] = history
        return jsonify({"reply": assistant_message})
    except Exception as e:
        return jsonify({"error": f"エラーが発生しました: {str(e)}"}), 500


@app.route("/company/<int:company_id>/delete", methods=["POST"])
def company_delete(company_id):
    delete_company(company_id)
    return redirect(url_for("index"))


# ─────────────────────────────────────────
# URL自動取得API（Work6フィールドに対応）
# ─────────────────────────────────────────

@app.route("/api/fetch_url", methods=["POST"])
def api_fetch_url():
    """HPのURLを読み込んでAIでWork6の会社情報を自動抽出"""
    data = request.get_json()
    url = data.get("url", "").strip()
    if not url:
        return jsonify({"error": "URLを入力してください"}), 400

    try:
        hp_text = fetch_url_text(url)
    except Exception as e:
        return jsonify({"error": f"URLの取得に失敗しました: {str(e)}"}), 400

    prompt = f"""以下は企業のHP（トップページ＋会社概要ページ）から取得したテキストです。
このテキストをもとに、下記の情報をJSON形式で正確に抽出してください。
情報が明確に記載されていない項目は空文字列にしてください。架空の情報は絶対に作らないでください。

抽出する項目（APS Work6 お客様概要）：
- company_name: 会社名（正式名称）
- industry: 業種（例：旅行業、製造業、IT・SIerなど）
- founded: 創業年（例：1980年）
- established: 設立年（例：1990年）
- headquarters: 本社所在地
- capital: 資本金
- employees: 従業員数（グループ全体含む場合はその旨も）
- revenue: 売上高（直近の数値）
- operating_profit: 経常利益
- branches: 事業所・店舗（主要なもの）
- group_companies: グループ会社・関連会社
- overview: 事業内容・会社概要（300字程度。具体的な事業・ビジネスモデル・特徴を記載）
- mvv: ミッション・ビジョン・バリュー（MVV）または企業理念・経営方針の概要
- president_profile: 代表者名と経歴（簡潔に150字程度）
- mid_term_plan: 中期経営計画・会社の方針（200字程度）

HPテキスト：
{hp_text}

JSON形式のみで返してください（説明文不要）："""

    try:
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=2000,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = response.content[0].text.strip()
        # JSON部分だけ抽出
        if "```" in raw:
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        result = json.loads(raw)
        return jsonify({"success": True, "data": result})
    except json.JSONDecodeError:
        return jsonify({"error": "AIの返答をパースできませんでした。もう一度試してください。"}), 500
    except Exception as e:
        return jsonify({"error": f"エラーが発生しました: {str(e)}"}), 500


# ─────────────────────────────────────────
# APS業界分析 自動生成API（Work2〜5・7）
# ─────────────────────────────────────────

@app.route("/company/<int:company_id>/generate_analysis", methods=["POST"])
def generate_analysis(company_id):
    """PEST・5Forces・SWOT・クロスSWOT・ポジショニング・投資エリアを自動生成してDBに保存"""
    company = get_company(company_id)
    if not company:
        return jsonify({"error": "会社情報が見つかりません"}), 404

    company_name = company.get("company_name", "")
    industry = company.get("industry", "")
    overview = company.get("overview", "")
    mid_term_plan = company.get("mid_term_plan", "")

    prompt = f"""以下の企業について、APS（アカウントプランニングセッション）の観点で業界分析を行ってください。
架空の情報は作らず、一般的な業界知識と提供された情報をもとに分析してください。

企業名：{company_name}
業種：{industry}
会社概要：{overview}
中期経営計画：{mid_term_plan}

以下の分析を日本語で行い、JSON形式で返してください：

{{
  "pest": "PEST分析（P:政治・法規制、E:経済・市場環境、S:社会・文化・人口動態、T:技術革新の4観点で業界への影響を具体的に分析。各観点2〜3点ずつ）",
  "five_forces": "5Forces分析（①新規参入の脅威、②売り手の交渉力、③買い手の交渉力、④代替品の脅威、⑤業界内競合の5観点と、最後に総合的な考察を記載）",
  "swot": "SWOT分析（強み：内部のプラス要因、弱み：内部のマイナス要因、機会：外部のプラス要因、脅威：外部のマイナス要因を各3〜5項目ずつ箇条書き）",
  "cross_swot": "クロスSWOT分析（強み×機会→積極戦略、強み×脅威→差別化戦略、弱み×機会→改善戦略、弱み×脅威→防衛戦略の4つの戦略方向性を記載）",
  "positioning": "ポジショニング（業界内での{company_name}の立ち位置、主要競合との差別化ポイント、競争優位性の源泉）",
  "investment_areas": "投資エリア（公開情報や業界トレンドから投資可能性がある領域を箇条書き。各項目に想定されるソリューション例も記載）"
}}

JSONのみ返してください（説明文不要）："""

    try:
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=3500,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = response.content[0].text.strip()
        if "```" in raw:
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        result = json.loads(raw)

        # DBに保存
        update_data = dict(company)
        update_data.update(result)
        update_company(company_id, update_data)

        return jsonify({"success": True, "data": result})
    except Exception as e:
        return jsonify({"error": f"分析生成に失敗しました: {str(e)}"}), 500


# ─────────────────────────────────────────
# PPTXダウンロード
# ─────────────────────────────────────────

@app.route("/company/<int:company_id>/download_ppt", methods=["GET"])
def download_ppt(company_id):
    """テンプレートに会社情報を流し込んでPPTXをダウンロード"""
    from pptx import Presentation

    company = get_company(company_id)
    if not company:
        return jsonify({"error": "会社情報が見つかりません"}), 404

    if not TEMPLATE_PATH.exists():
        return jsonify({"error": "テンプレートファイルが見つかりません"}), 500

    prs = Presentation(str(TEMPLATE_PATH))
    slides = prs.slides

    company_name = company.get("company_name", "")
    sales_person = company.get("sales_person", "")

    def safe_write(slide, keyword, text, max_chars=800):
        """キーワードを含むシェイプにテキストを書き込む"""
        if not text:
            return
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if keyword in shape.text_frame.text:
                tf = shape.text_frame
                tf.word_wrap = True
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = ""
                if tf.paragraphs:
                    if tf.paragraphs[0].runs:
                        tf.paragraphs[0].runs[0].text = str(text)[:max_chars]
                    else:
                        tf.paragraphs[0].text = str(text)[:max_chars]
                return

    def replace_in_slide(slide, old, new):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)

    # ── P1: タイトル ──
    if len(slides) > 0:
        replace_in_slide(slides[0], "N.K.ナーツ株式会社", f"{company_name}\n担当：{sales_person}")

    # ── P13 (Work6: お客様概要) ──
    if len(slides) > 12:
        slide = slides[12]
        overview_text = f"""社名：{company_name}
創業：{company.get('founded', '')}　設立：{company.get('established', '')}
代表者：{company.get('president_profile', '').split('\\n')[0] if company.get('president_profile') else ''}
本社：{company.get('headquarters', '')}
資本金：{company.get('capital', '')}　売上高：{company.get('revenue', '')}
経常利益：{company.get('operating_profit', '')}　社員数：{company.get('employees', '')}

【事業内容】
{company.get('overview', '')}

【グループ会社】
{company.get('group_companies', '')}"""
        safe_write(slide, "社名", overview_text[:900])

    # ── P18 (Work7: MVV) ──
    if len(slides) > 17:
        slide = slides[17]
        mvv_text = company.get("mvv", "")
        if mvv_text:
            safe_write(slide, "MVV", mvv_text[:600])

    # ── P20 (Work7: 経営方針・中計) ──
    if len(slides) > 19:
        slide = slides[19]
        plan_text = f"""【中期経営計画・経営方針】
{company.get('mid_term_plan', '')}

【IR情報】
{company.get('ir_info', '')}"""
        safe_write(slide, "経営方針", plan_text[:800])

    # ── P22 (Work7: 投資エリア) ──
    if len(slides) > 21:
        slide = slides[21]
        invest_text = company.get("investment_areas", "")
        if invest_text:
            # 長いテキストボックスを探す
            for shape in slide.shapes:
                if shape.has_text_frame and len(shape.text_frame.text) > 50:
                    tf = shape.text_frame
                    tf.word_wrap = True
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.text = ""
                    if tf.paragraphs:
                        if tf.paragraphs[0].runs:
                            tf.paragraphs[0].runs[0].text = invest_text[:800]
                        else:
                            tf.paragraphs[0].text = invest_text[:800]
                    break

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename = f"APS_{company_name}_{sales_person}.pptx"
    filename = filename.replace("/", "_").replace(" ", "_")

    return send_file(
        buf,
        as_attachment=True,
        download_name=filename,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


# ─────────────────────────────────────────
# フォームパース
# ─────────────────────────────────────────

def _parse_form(form):
    system_names = form.getlist("system_name[]")
    system_types = form.getlist("system_type[]")
    system_installed = form.getlist("system_installed[]")
    system_replace = form.getlist("system_replace[]")
    systems = []
    for i in range(len(system_names)):
        if system_names[i].strip():
            systems.append({
                "name": system_names[i].strip(),
                "type": system_types[i].strip() if i < len(system_types) else "",
                "installed": system_installed[i].strip() if i < len(system_installed) else "",
                "replace": system_replace[i].strip() if i < len(system_replace) else "",
            })

    person_names = form.getlist("person_name[]")
    person_titles = form.getlist("person_title[]")
    person_decisions = form.getlist("person_decision[]")
    person_notes = form.getlist("person_note[]")
    key_persons = []
    for i in range(len(person_names)):
        if person_names[i].strip():
            key_persons.append({
                "name": person_names[i].strip(),
                "title": person_titles[i].strip() if i < len(person_titles) else "",
                "decision": person_decisions[i].strip() if i < len(person_decisions) else "不明",
                "note": person_notes[i].strip() if i < len(person_notes) else "",
            })

    return {
        "company_name": form.get("company_name", "").strip(),
        "industry": form.get("industry", "").strip(),
        "employees": form.get("employees", "").strip(),
        "revenue": form.get("revenue", "").strip(),
        "hp_url": form.get("hp_url", "").strip(),
        "sales_person": form.get("sales_person", "").strip(),
        # Work6フィールド
        "founded": form.get("founded", "").strip(),
        "established": form.get("established", "").strip(),
        "headquarters": form.get("headquarters", "").strip(),
        "capital": form.get("capital", "").strip(),
        "operating_profit": form.get("operating_profit", "").strip(),
        "branches": form.get("branches", "").strip(),
        "group_companies": form.get("group_companies", "").strip(),
        "company_detail": form.get("company_detail", "").strip(),
        "overview": form.get("overview", "").strip(),
        "president_profile": form.get("president_profile", "").strip(),
        "mvv": form.get("mvv", "").strip(),
        # Work7フィールド
        "mid_term_plan": form.get("mid_term_plan", "").strip(),
        "ir_info": form.get("ir_info", "").strip(),
        "investment_areas": form.get("investment_areas", "").strip(),
        # 業界分析
        "pest": form.get("pest", "").strip(),
        "five_forces": form.get("five_forces", "").strip(),
        "swot": form.get("swot", "").strip(),
        "cross_swot": form.get("cross_swot", "").strip(),
        "positioning": form.get("positioning", "").strip(),
        # 営業調査
        "systems": systems,
        "key_persons": key_persons,
        "competitors": form.get("competitors", "").strip(),
        "end_user_issues": form.get("end_user_issues", "").strip(),
        "latent_needs": form.get("latent_needs", "").strip(),
        "big_play": form.get("big_play", "").strip(),
        "pipeline": form.get("pipeline", "").strip(),
        # Work8〜13
        "activity_history": form.get("activity_history", "").strip(),
        "mid_long_term_plan": form.get("mid_long_term_plan", "").strip(),
        "org_chart": form.get("org_chart", "").strip(),
        "forecast": form.get("forecast", "").strip(),
        "key_cases": form.get("key_cases", "").strip(),
        "coverage_map": form.get("coverage_map", "").strip(),
        "action_plan": form.get("action_plan", "").strip(),
        "company_requests": form.get("company_requests", "").strip(),
    }


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5002))
    app.run(debug=False, host="0.0.0.0", port=port)
